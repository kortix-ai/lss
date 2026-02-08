"""Tests for lss_extract — document format extraction layer.

TDD: These tests are written BEFORE the implementation.
Each test creates a real file of the target format and verifies extraction.
"""

import email
import email.mime.text
import email.mime.multipart
from pathlib import Path

import pytest


# ── PDF extraction ───────────────────────────────────────────────────────────


def _make_pdf(path, text="Hello World from PDF"):
    """Create a minimal PDF with extractable text using pdfminer-compatible format."""
    # Use reportlab if available, otherwise write a minimal raw PDF
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(path), pagesize=letter)
        c.drawString(72, 720, text)
        c.save()
    except ImportError:
        # Minimal raw PDF with text
        content = (
            "%PDF-1.4\n"
            "1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            "2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            "3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
            "/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
            "4 0 obj<</Length 44>>stream\n"
            "BT /F1 12 Tf 72 720 Td (" + text + ") Tj ET\n"
            "endstream\nendobj\n"
            "5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
            "xref\n0 6\n"
            "0000000000 65535 f \n"
            "0000000009 00000 n \n"
            "0000000058 00000 n \n"
            "0000000115 00000 n \n"
            "0000000266 00000 n \n"
            "0000000360 00000 n \n"
            "trailer<</Size 6/Root 1 0 R>>\n"
            "startxref\n431\n%%EOF"
        )
        path.write_text(content)


class TestExtractPdf:
    def test_basic(self, tmp_path):
        from lss_extract import extract_pdf

        pdf = tmp_path / "test.pdf"
        _make_pdf(pdf, "Kubernetes deployment pipeline")
        text = extract_pdf(str(pdf))
        assert "Kubernetes" in text or "deployment" in text or "pipeline" in text

    def test_empty_pdf(self, tmp_path):
        """Corrupt/empty PDF returns empty string, not an exception."""
        from lss_extract import extract_pdf

        pdf = tmp_path / "empty.pdf"
        pdf.write_bytes(b"%PDF-1.4\n%%EOF")
        text = extract_pdf(str(pdf))
        assert isinstance(text, str)

    def test_nonexistent_pdf(self, tmp_path):
        """Missing file returns empty string."""
        from lss_extract import extract_pdf

        text = extract_pdf(str(tmp_path / "missing.pdf"))
        assert text == ""


# ── DOCX extraction ─────────────────────────────────────────────────────────


def _make_docx(path, paragraphs=None, headings=None, table_data=None):
    """Create a DOCX file with optional paragraphs, headings, and table."""
    from docx import Document

    doc = Document()
    if headings:
        for level, text in headings:
            doc.add_heading(text, level=level)
    if paragraphs:
        for p in paragraphs:
            doc.add_paragraph(p)
    if table_data:
        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
        for i, row in enumerate(table_data):
            for j, cell in enumerate(row):
                table.rows[i].cells[j].text = str(cell)
    doc.save(str(path))


class TestExtractDocx:
    def test_paragraphs(self, tmp_path):
        from lss_extract import extract_docx

        docx = tmp_path / "test.docx"
        _make_docx(docx, paragraphs=["Authentication uses JWT tokens", "Database uses PostgreSQL"])
        text = extract_docx(str(docx))
        assert "JWT tokens" in text
        assert "PostgreSQL" in text

    def test_headings(self, tmp_path):
        from lss_extract import extract_docx

        docx = tmp_path / "test.docx"
        _make_docx(
            docx,
            headings=[(1, "Architecture Overview"), (2, "Database Layer")],
            paragraphs=["The system uses microservices"],
        )
        text = extract_docx(str(docx))
        assert "Architecture Overview" in text
        assert "Database Layer" in text
        assert "microservices" in text

    def test_tables(self, tmp_path):
        from lss_extract import extract_docx

        docx = tmp_path / "test.docx"
        _make_docx(docx, table_data=[["Name", "Role"], ["Alice", "Engineer"], ["Bob", "Designer"]])
        text = extract_docx(str(docx))
        assert "Alice" in text
        assert "Engineer" in text
        assert "Bob" in text

    def test_empty_docx(self, tmp_path):
        from lss_extract import extract_docx

        docx = tmp_path / "empty.docx"
        _make_docx(docx)
        text = extract_docx(str(docx))
        assert isinstance(text, str)

    def test_nonexistent_docx(self, tmp_path):
        from lss_extract import extract_docx

        text = extract_docx(str(tmp_path / "missing.docx"))
        assert text == ""


# ── XLSX extraction ─────────────────────────────────────────────────────────


def _make_xlsx(path, sheets=None):
    """Create an XLSX file with named sheets and cell data.
    sheets: dict of {sheet_name: [[row1_data], [row2_data], ...]}
    """
    from openpyxl import Workbook

    wb = Workbook()
    if sheets:
        first = True
        for name, rows in sheets.items():
            if first:
                ws = wb.active
                ws.title = name
                first = False
            else:
                ws = wb.create_sheet(name)
            for row in rows:
                ws.append(row)
    wb.save(str(path))


class TestExtractXlsx:
    def test_basic(self, tmp_path):
        from lss_extract import extract_xlsx

        xlsx = tmp_path / "test.xlsx"
        _make_xlsx(xlsx, {"Data": [["Name", "Score"], ["Alice", 95], ["Bob", 87]]})
        text = extract_xlsx(str(xlsx))
        assert "Alice" in text
        assert "95" in text
        assert "Bob" in text

    def test_multiple_sheets(self, tmp_path):
        from lss_extract import extract_xlsx

        xlsx = tmp_path / "multi.xlsx"
        _make_xlsx(
            xlsx,
            {
                "Users": [["Name"], ["Charlie"]],
                "Products": [["Item"], ["Widget"]],
            },
        )
        text = extract_xlsx(str(xlsx))
        assert "Charlie" in text
        assert "Widget" in text
        # Sheet names should be included
        assert "Users" in text
        assert "Products" in text

    def test_empty_xlsx(self, tmp_path):
        from lss_extract import extract_xlsx

        xlsx = tmp_path / "empty.xlsx"
        _make_xlsx(xlsx)
        text = extract_xlsx(str(xlsx))
        assert isinstance(text, str)

    def test_nonexistent_xlsx(self, tmp_path):
        from lss_extract import extract_xlsx

        text = extract_xlsx(str(tmp_path / "missing.xlsx"))
        assert text == ""


# ── PPTX extraction ─────────────────────────────────────────────────────────


def _make_pptx(path, slides=None):
    """Create a PPTX file with slides.
    slides: list of (title, body, notes) tuples.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    if slides:
        for title, body, notes in slides:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))


class TestExtractPptx:
    def test_slides(self, tmp_path):
        from lss_extract import extract_pptx

        pptx = tmp_path / "test.pptx"
        _make_pptx(
            pptx,
            [
                ("Architecture", "Microservices with Kubernetes", "Speaker notes here"),
                ("Database", "PostgreSQL with read replicas", None),
            ],
        )
        text = extract_pptx(str(pptx))
        assert "Architecture" in text
        assert "Microservices" in text
        assert "PostgreSQL" in text

    def test_speaker_notes(self, tmp_path):
        from lss_extract import extract_pptx

        pptx = tmp_path / "notes.pptx"
        _make_pptx(pptx, [("Title", "Body", "Important deployment notes")])
        text = extract_pptx(str(pptx))
        assert "deployment notes" in text

    def test_empty_pptx(self, tmp_path):
        from lss_extract import extract_pptx

        pptx = tmp_path / "empty.pptx"
        from pptx import Presentation

        Presentation().save(str(pptx))
        text = extract_pptx(str(pptx))
        assert isinstance(text, str)

    def test_nonexistent_pptx(self, tmp_path):
        from lss_extract import extract_pptx

        text = extract_pptx(str(tmp_path / "missing.pptx"))
        assert text == ""


# ── HTML extraction ──────────────────────────────────────────────────────────


class TestExtractHtml:
    def test_strips_tags(self, tmp_path):
        from lss_extract import extract_html

        html = tmp_path / "test.html"
        html.write_text("<html><body><p>Hello World</p><p>Second paragraph</p></body></html>")
        text = extract_html(str(html))
        assert "Hello World" in text
        assert "Second paragraph" in text
        assert "<p>" not in text
        assert "<html>" not in text

    def test_skips_script_and_style(self, tmp_path):
        from lss_extract import extract_html

        html = tmp_path / "script.html"
        html.write_text(
            "<html><head><style>body{color:red}</style></head>"
            "<body><script>alert('xss')</script>"
            "<p>Actual content here</p></body></html>"
        )
        text = extract_html(str(html))
        assert "Actual content" in text
        assert "alert" not in text
        assert "color:red" not in text

    def test_extracts_headings(self, tmp_path):
        from lss_extract import extract_html

        html = tmp_path / "headings.html"
        html.write_text(
            "<html><body>"
            "<h1>Main Title</h1>"
            "<h2>Section One</h2>"
            "<p>Content of section one</p>"
            "</body></html>"
        )
        text = extract_html(str(html))
        assert "Main Title" in text
        assert "Section One" in text
        assert "Content of section one" in text

    def test_extracts_links_text(self, tmp_path):
        from lss_extract import extract_html

        html = tmp_path / "links.html"
        html.write_text('<html><body><a href="https://example.com">Click here</a></body></html>')
        text = extract_html(str(html))
        assert "Click here" in text

    def test_empty_html(self, tmp_path):
        from lss_extract import extract_html

        html = tmp_path / "empty.html"
        html.write_text("")
        text = extract_html(str(html))
        assert isinstance(text, str)

    def test_nonexistent_html(self, tmp_path):
        from lss_extract import extract_html

        text = extract_html(str(tmp_path / "missing.html"))
        assert text == ""


# ── Email extraction ─────────────────────────────────────────────────────────


def _make_eml(path, subject="Test Email", body="This is the email body", html_body=None):
    """Create a .eml file."""
    if html_body:
        msg = email.mime.multipart.MIMEMultipart("alternative")
        msg.attach(email.mime.text.MIMEText(body, "plain"))
        msg.attach(email.mime.text.MIMEText(html_body, "html"))
    else:
        msg = email.mime.text.MIMEText(body)
    msg["Subject"] = subject
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    path.write_text(msg.as_string())


class TestExtractEmail:
    def test_body(self, tmp_path):
        from lss_extract import extract_email

        eml = tmp_path / "test.eml"
        _make_eml(eml, body="Meeting notes from the architecture review")
        text = extract_email(str(eml))
        assert "architecture review" in text

    def test_subject(self, tmp_path):
        from lss_extract import extract_email

        eml = tmp_path / "test.eml"
        _make_eml(eml, subject="Q4 Planning Session", body="Discussion of roadmap")
        text = extract_email(str(eml))
        assert "Q4 Planning Session" in text
        assert "roadmap" in text

    def test_multipart(self, tmp_path):
        from lss_extract import extract_email

        eml = tmp_path / "multi.eml"
        _make_eml(
            eml,
            subject="HTML Email",
            body="Plain text version of email content",
            html_body="<html><body><p>HTML version of email content</p></body></html>",
        )
        text = extract_email(str(eml))
        # Should extract at least one version
        assert "email content" in text

    def test_no_mime_headers_in_output(self, tmp_path):
        from lss_extract import extract_email

        eml = tmp_path / "test.eml"
        _make_eml(eml, body="Just the body text")
        text = extract_email(str(eml))
        assert "Content-Type" not in text
        assert "MIME-Version" not in text

    def test_empty_email(self, tmp_path):
        from lss_extract import extract_email

        eml = tmp_path / "empty.eml"
        eml.write_text("")
        text = extract_email(str(eml))
        assert isinstance(text, str)

    def test_nonexistent_email(self, tmp_path):
        from lss_extract import extract_email

        text = extract_email(str(tmp_path / "missing.eml"))
        assert text == ""


# ── Dispatcher (extract_text) ────────────────────────────────────────────────


class TestExtractDispatcher:
    """Test the main extract_text() dispatcher routes to the right extractor."""

    def test_routes_pdf(self, tmp_path):
        from lss_extract import extract_text

        pdf = tmp_path / "test.pdf"
        _make_pdf(pdf, "dispatcher test PDF")
        text = extract_text(str(pdf))
        assert isinstance(text, str)

    def test_routes_docx(self, tmp_path):
        from lss_extract import extract_text

        docx = tmp_path / "test.docx"
        _make_docx(docx, paragraphs=["dispatcher test DOCX"])
        text = extract_text(str(docx))
        assert "dispatcher test DOCX" in text

    def test_routes_xlsx(self, tmp_path):
        from lss_extract import extract_text

        xlsx = tmp_path / "test.xlsx"
        _make_xlsx(xlsx, {"Sheet1": [["dispatcher", "test"]]})
        text = extract_text(str(xlsx))
        assert "dispatcher" in text

    def test_routes_pptx(self, tmp_path):
        from lss_extract import extract_text

        pptx = tmp_path / "test.pptx"
        _make_pptx(pptx, [("Dispatcher", "test PPTX content", None)])
        text = extract_text(str(pptx))
        assert "Dispatcher" in text

    def test_routes_html(self, tmp_path):
        from lss_extract import extract_text

        html = tmp_path / "test.html"
        html.write_text("<p>dispatcher test HTML</p>")
        text = extract_text(str(html))
        assert "dispatcher test HTML" in text

    def test_routes_eml(self, tmp_path):
        from lss_extract import extract_text

        eml = tmp_path / "test.eml"
        _make_eml(eml, body="dispatcher test EML")
        text = extract_text(str(eml))
        assert "dispatcher test EML" in text

    def test_routes_json(self, tmp_path):
        from lss_extract import extract_text

        f = tmp_path / "test.json"
        f.write_text('{"key": "dispatcher test JSON"}')
        text = extract_text(str(f))
        assert "dispatcher test JSON" in text

    def test_routes_csv(self, tmp_path):
        from lss_extract import extract_text

        f = tmp_path / "test.csv"
        f.write_text("name,value\nAlice,100\n")
        text = extract_text(str(f))
        assert "Alice" in text

    def test_routes_plain_text(self, tmp_path):
        from lss_extract import extract_text

        f = tmp_path / "test.txt"
        f.write_text("plain text content")
        text = extract_text(str(f))
        assert "plain text content" in text

    def test_routes_python(self, tmp_path):
        from lss_extract import extract_text

        f = tmp_path / "test.py"
        f.write_text("def hello():\n    return 'world'\n")
        text = extract_text(str(f))
        assert "hello" in text

    def test_routes_markdown(self, tmp_path):
        from lss_extract import extract_text

        f = tmp_path / "test.md"
        f.write_text("# Title\nMarkdown content here\n")
        text = extract_text(str(f))
        assert "Markdown content" in text
