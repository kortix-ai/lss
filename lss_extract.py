"""lss_extract — Document format extraction layer.

Converts files of various formats into plain text for indexing.
Each extractor is self-contained and handles its own errors gracefully
(returns "" on failure rather than raising).

Supported formats:
  PDF   — pdfminer.six (layout-aware, handles multi-column)
  DOCX  — python-docx (paragraphs, headings, tables)
  XLSX  — openpyxl (all sheets, cell values)
  PPTX  — python-pptx (slides, speaker notes)
  HTML  — beautifulsoup4 (strips tags, skips script/style)
  EML   — stdlib email (subject + body, handles multipart)
  JSON  — stdlib json (flattens values)
  JSONL — stdlib json (line-by-line)
  CSV   — stdlib csv (flattens rows)
  Text  — stdlib (utf-8 with fallback encodings)
"""

import json
import re
from pathlib import Path


# ── PDF ──────────────────────────────────────────────────────────────────────


def extract_pdf(file_path: str) -> str:
    """Extract text from a PDF using pdfminer.six (layout-aware)."""
    try:
        from pdfminer.high_level import extract_text as _pdfminer_extract
        text = _pdfminer_extract(file_path)
        return text.strip() if text else ""
    except ImportError:
        # Fallback to PyPDF2 if pdfminer not installed
        try:
            import PyPDF2
            text = ""
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.strip()
        except Exception:
            return ""
    except Exception:
        return ""


# ── DOCX ─────────────────────────────────────────────────────────────────────


def extract_docx(file_path: str) -> str:
    """Extract text from a DOCX file (paragraphs, headings, tables)."""
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []

        # Paragraphs (includes headings — they're paragraphs with heading styles)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        return "\n".join(parts)
    except Exception:
        return ""


# ── XLSX ─────────────────────────────────────────────────────────────────────


def extract_xlsx(file_path: str) -> str:
    """Extract text from an XLSX file (all sheets, cell values)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[{sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                row_texts = [str(cell) for cell in row if cell is not None]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        wb.close()
        return "\n".join(parts)
    except Exception:
        return ""


# ── PPTX ─────────────────────────────────────────────────────────────────────


def extract_pptx(file_path: str) -> str:
    """Extract text from a PPTX file (slide text + speaker notes)."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                # Tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(notes)

            if slide_texts:
                parts.append("\n".join(slide_texts))

        return "\n\n".join(parts)
    except Exception:
        return ""


# ── HTML ─────────────────────────────────────────────────────────────────────


def extract_html(file_path: str) -> str:
    """Extract text from HTML, stripping tags and skipping script/style."""
    try:
        text = Path(file_path).read_text(encoding="utf-8", errors="replace")
        if not text.strip():
            return ""

        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")

        # Remove script and style elements
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()

        # Get text with newlines between block elements
        clean = soup.get_text(separator="\n", strip=True)
        # Collapse multiple blank lines
        clean = re.sub(r"\n{3,}", "\n\n", clean)
        return clean.strip()
    except Exception:
        return ""


# ── Email ────────────────────────────────────────────────────────────────────


def extract_email(file_path: str) -> str:
    """Extract subject + body from an .eml file."""
    try:
        import email as email_mod
        from email.policy import default as default_policy

        raw = Path(file_path).read_text(encoding="utf-8", errors="replace")
        if not raw.strip():
            return ""

        msg = email_mod.message_from_string(raw, policy=default_policy)
        parts = []

        # Subject
        subject = msg.get("Subject", "")
        if subject:
            parts.append(f"Subject: {subject}")

        # Body — walk multipart, prefer plain text
        body_text = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    payload = part.get_content()
                    if isinstance(payload, str):
                        body_text = payload
                        break
            # Fallback to HTML part if no plain text
            if not body_text:
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct == "text/html":
                        payload = part.get_content()
                        if isinstance(payload, str):
                            # Strip HTML tags
                            try:
                                from bs4 import BeautifulSoup
                                soup = BeautifulSoup(payload, "html.parser")
                                body_text = soup.get_text(separator="\n", strip=True)
                            except ImportError:
                                body_text = re.sub(r"<[^>]+>", " ", payload)
                        break
        else:
            payload = msg.get_content()
            if isinstance(payload, str):
                body_text = payload

        if body_text:
            parts.append(body_text.strip())

        return "\n\n".join(parts)
    except Exception:
        return ""


# ── JSON ─────────────────────────────────────────────────────────────────────


def extract_json(file_path: str) -> str:
    """Extract text from a JSON file by flattening string values."""
    try:
        text = Path(file_path).read_text(encoding="utf-8")
        data = json.loads(text)
        if isinstance(data, dict):
            return " ".join(str(v) for v in data.values() if isinstance(v, (str, int, float)))
        elif isinstance(data, list):
            return " ".join(str(item) for item in data if isinstance(item, (str, int, float)))
        else:
            return str(data)
    except Exception:
        return ""


# ── JSONL ────────────────────────────────────────────────────────────────────


def extract_jsonl(file_path: str) -> str:
    """Extract text from a JSONL file (one JSON object per line)."""
    try:
        text = Path(file_path).read_text(encoding="utf-8")
        texts = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            obj = json.loads(line)
            if isinstance(obj, dict):
                if "text" in obj:
                    texts.append(obj["text"])
                elif "title" in obj and "text" in obj:
                    texts.append(f"{obj['title']} {obj['text']}")
                else:
                    texts.append(
                        " ".join(str(v) for v in obj.values() if isinstance(v, (str, int, float)))
                    )
        return "\n".join(texts)
    except Exception:
        return ""


# ── CSV ──────────────────────────────────────────────────────────────────────


def extract_csv(file_path: str) -> str:
    """Extract text from a CSV file by flattening rows."""
    try:
        import csv
        import io

        text = Path(file_path).read_text(encoding="utf-8")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return ""
        result = " ".join(rows[0])  # header
        for row in rows[1 : min(100, len(rows))]:
            result += " " + " ".join(row)
        return result
    except Exception:
        return ""


# ── Plain text ───────────────────────────────────────────────────────────────


def extract_plain_text(file_path: str) -> str:
    """Read a plain text file with encoding fallbacks."""
    path = Path(file_path)
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        for enc in ("latin-1", "cp1252", "iso-8859-1"):
            try:
                return path.read_text(encoding=enc)
            except UnicodeDecodeError:
                continue
    except Exception:
        pass
    return ""


# ── Dispatcher ───────────────────────────────────────────────────────────────

# Extension → extractor function mapping
_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,  # openpyxl can handle some .xls via compatibility
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".eml": extract_email,
    ".json": extract_json,
    ".jsonl": extract_jsonl,
    ".csv": extract_csv,
    ".tsv": extract_csv,  # csv.reader handles TSV with delimiter detection
}


def extract_text(file_path: str) -> str:
    """Extract text from any supported file format.

    Routes to format-specific extractors based on file extension.
    Falls back to plain text reading for unknown extensions.

    Returns empty string on any failure (never raises).
    """
    path = Path(file_path)
    if not path.exists():
        return ""

    ext = path.suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor:
        return extractor(file_path)

    # Default: plain text
    return extract_plain_text(file_path)
